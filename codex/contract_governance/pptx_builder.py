"""PPTX renderer for the ABACUS Contract Governance Workbench."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .io import content_hash
from .schema import GovernanceSSOT

GENERATOR = "CODEX Contract Governance Generator"
MAX_ROWS_PER_SLIDE = 15
_TITLE_LAYOUT = 0
_TITLE_ONLY_LAYOUT = 5


def build_pptx(payload: dict[str, object], ssot: GovernanceSSOT, path: Path) -> None:
    """Render a governance payload as a tiered PPTX deck."""

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    title_slide.shapes.title.text = (
        f"{payload['package_id']} {str(payload['tier']).upper()} Governance Snapshot"
    )
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Content hash (sha256): {content_hash(payload)}"

    for sheet_payload in payload["sheets"]:  # type: ignore[index]
        columns = sheet_payload["columns"]
        rows = sheet_payload["rows"]
        for chunk_index, chunk in enumerate(_chunk_rows(rows, MAX_ROWS_PER_SLIDE)):
            heading = sheet_payload["name"]
            if chunk_index > 0:
                heading = f"{heading} (cont.)"
            _add_table_slide(prs, heading, columns, chunk)

    fixed = ssot.build.fixed_docprops_timestamp.replace(tzinfo=None)
    prs.core_properties.created = fixed
    prs.core_properties.modified = fixed
    prs.core_properties.author = GENERATOR
    prs.core_properties.last_modified_by = GENERATOR

    prs.save(path)


def _chunk_rows(
    rows: list[dict[str, str]], size: int
) -> list[list[dict[str, str]]]:
    if not rows:
        return [[]]
    return [rows[index : index + size] for index in range(0, len(rows), size)]


def _add_table_slide(
    prs: Presentation,
    heading: str,
    columns: list[str],
    rows: list[dict[str, str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_ONLY_LAYOUT])
    slide.shapes.title.text = heading

    table_rows = 1 + len(rows)
    table_cols = len(columns)
    left, top = Inches(0.5), Inches(1.5)
    width, height = Inches(9.0), Inches(0.5) * table_rows
    graphic_frame = slide.shapes.add_table(table_rows, table_cols, left, top, width, height)
    table = graphic_frame.table

    for col_index, column in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = column
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)

    for row_index, row in enumerate(rows, start=1):
        for col_index, column in enumerate(columns):
            cell = table.cell(row_index, col_index)
            cell.text = str(row.get(column, ""))
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
