from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation

from codex.contract_governance.builder import build_artifacts
from codex.contract_governance.io import load_ssot
from codex.contract_governance.validator import validate_generated

SSOT = Path("contract_governance/ssot/abacus_contract_governance.yaml")
FIXED_DOC_PROPS = datetime(2026, 1, 1, 0, 0, 0)
GENERATOR = "CODEX Contract Governance Generator"
FORBIDDEN_BIDDER_STRINGS = {
    "Extraction Audit",
    "Evaluation Notes",
    "REQ-W001-001",
    "Generated tier stripping",
    "Internal evaluation placeholder",
}


def _pptx_full_text(path: Path) -> str:
    prs = Presentation(path)
    fragments: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                fragments.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fragments.append(cell.text)
    return "\n".join(fragments)


def test_build_artifacts_generates_a_pptx_deck(tmp_path: Path) -> None:
    ssot = load_ssot(SSOT)
    result = build_artifacts(ssot, tmp_path, "internal")

    assert "pptx" in result
    pptx_path = Path(result["pptx"])
    assert pptx_path.exists()

    deck_text = _pptx_full_text(pptx_path)
    assert "Requirements" in deck_text
    assert "Extraction Audit" in deck_text
    assert "REQ-W001-001" in deck_text


def test_pptx_pinned_core_properties_are_deterministic(tmp_path: Path) -> None:
    ssot = load_ssot(SSOT)
    first = build_artifacts(ssot, tmp_path / "first", "internal")
    second = build_artifacts(ssot, tmp_path / "second", "internal")

    for pptx_path in (first["pptx"], second["pptx"]):
        prs = Presentation(pptx_path)
        assert prs.core_properties.created.replace(tzinfo=None) == FIXED_DOC_PROPS
        assert prs.core_properties.modified.replace(tzinfo=None) == FIXED_DOC_PROPS
        assert prs.core_properties.author == GENERATOR
        assert prs.core_properties.last_modified_by == GENERATOR


def test_bidder_pptx_strips_internal_only_content(tmp_path: Path) -> None:
    ssot = load_ssot(SSOT)
    result = build_artifacts(ssot, tmp_path, "bidder")
    validate_generated(ssot, tmp_path, "bidder")

    pptx_text = _pptx_full_text(Path(result["pptx"]))
    for forbidden in FORBIDDEN_BIDDER_STRINGS:
        assert forbidden not in pptx_text
