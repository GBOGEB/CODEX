"""PPTX → digital-twin slide preview generator.

Ported from GBOGEB/CODESPACES_jyperter (tools/extract_slide_previews.py) and
adapted to CODEX conventions (main() guard, Path.mkdir, html.escape).

Produces three artifacts in *outdir*:

- ``slide_texts.json``       — full per-slide text content
- ``slide_digital_twins.json`` — compact summaries + styling hints
- ``slide_preview.html``     — browser-viewable HTML card deck

Usage::

    python src/tools/slide_preview_generator.py \\
        --pptx path/to/deck.pptx \\
        --outdir outputs/html/slide_previews \\
        --limit 20
"""

from __future__ import annotations

import argparse
import html
import json
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:
    def Presentation(*args: Any, **kwargs: Any) -> Any:
        raise ImportError(
            "python-pptx is required to run this tool.\n"
            "Install it with:  pip install python-pptx"
        )


# ── CSS embedded into the preview HTML ─────────────────────────────────────

_PREVIEW_CSS = """
<style>
  body {
    font-family: Arial, Helvetica, sans-serif;
    background: #111;
    color: #eee;
    margin: 0;
    padding: 24px;
  }
  .deck { display: flex; flex-wrap: wrap; gap: 16px; }
  .card {
    width: 360px;
    min-height: 200px;
    border-radius: 10px;
    box-shadow: 0 2px 12px rgba(0,0,0,.4);
    padding: 12px;
    overflow: auto;
  }
  .title { margin: 4px 0 6px 0; color: #111; }
  .meta { font-size: 12px; color: #444; margin-bottom: 6px; }
  .para { margin: 4px 0; color: #222; font-size: 12px; }
  .num {
    background: #444;
    color: #fff;
    border-radius: 6px;
    padding: 2px 6px;
    font-size: 11px;
  }
</style>
"""


# ── Core extraction ─────────────────────────────────────────────────────────

def extract_slides(pptx_path: Path, limit: int | None = None) -> list[dict[str, Any]]:
    """Extract text, shape metadata, and style hints from a PPTX file.

    Parameters
    ----------
    pptx_path:
        Path to the ``.pptx`` file.
    limit:
        Maximum number of slides to process.  ``None`` processes all slides.

    Returns
    -------
    list[dict]
        One dict per slide with keys: ``n``, ``texts``, ``shape_count``,
        ``bullets``, ``style_hint``.
    """
    prs = Presentation(str(pptx_path))
    all_slides = list(prs.slides)
    if limit is not None:
        all_slides = all_slides[:limit]

    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(all_slides):
        texts: list[str] = []
        bullets = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
                    if any(text.lstrip().startswith(ch) for ch in ("-", "•", "*")):
                        bullets += 1

        title = texts[0] if texts else ""
        is_title_slide = bool(title) and len(title) < 120 and title.isupper()
        if is_title_slide:
            bg = "#f0f8ff"
        elif bullets > 0:
            bg = "#ffffff"
        else:
            bg = "#fafafa"

        slides.append(
            {
                "n": i + 1,
                "texts": texts,
                "shape_count": len(slide.shapes),
                "bullets": bullets,
                "style_hint": {"bg": bg, "is_title": is_title_slide},
            }
        )
    return slides


def build_digital_twins(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Compact summary representation for each slide (digital twin format)."""
    return [
        {
            "slide": s["n"],
            "summary": s["texts"][:3],
            "style_hint": s["style_hint"],
            "shape_count": s["shape_count"],
            "bullets": s["bullets"],
        }
        for s in slides
    ]


def write_html_preview(slides: list[dict[str, Any]], out_html: Path) -> None:
    """Write a self-contained HTML card deck preview for the extracted slides.

    All user-supplied text is escaped via ``html.escape`` before insertion.
    """
    parts = [
        "<html>",
        "<head><meta charset='utf-8'><title>Slide Preview</title>",
        _PREVIEW_CSS,
        "</head><body>",
        f"<h2>Slide Preview ({len(slides)} slides)</h2>",
        "<div class='deck'>",
    ]
    for s in slides:
        bg = s.get("style_hint", {}).get("bg", "#ffffff")
        texts = s.get("texts", [])
        title = texts[0] if texts else ""
        parts.append(f"<div class='card' style='background:{bg}'>")
        parts.append(
            f"<div class='meta'>"
            f"<span class='num'>#{s['n']}</span>"
            f" &bull; shapes:{s['shape_count']}"
            f" &bull; bullets:{s['bullets']}"
            f"</div>"
        )
        if title:
            parts.append(f"<h3 class='title'>{html.escape(title[:180])}</h3>")
        for text in texts[1:5]:
            parts.append(f"<div class='para'>{html.escape(text[:220])}</div>")
        parts.append("</div>")
    parts += ["</div>", "</body>", "</html>"]
    out_html.write_text("\n".join(parts), encoding="utf-8")


# ── CLI entry point ─────────────────────────────────────────────────────────

def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Generate slide text + digital twin previews from a PPTX file."
    )
    parser.add_argument("--pptx", required=True, help="Path to source .pptx file.")
    parser.add_argument(
        "--outdir",
        required=True,
        help="Output directory for generated artifacts.",
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help="Maximum number of slides to process (default: all).",
    )
    args = parser.parse_args(argv)

    pptx_path = Path(args.pptx)
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    slides = extract_slides(pptx_path, args.limit)
    twins = build_digital_twins(slides)

    texts_out = outdir / "slide_texts.json"
    twins_out = outdir / "slide_digital_twins.json"
    html_out = outdir / "slide_preview.html"

    texts_out.write_text(json.dumps(slides, indent=2, ensure_ascii=False), encoding="utf-8")
    twins_out.write_text(json.dumps(twins, indent=2, ensure_ascii=False), encoding="utf-8")
    write_html_preview(slides, html_out)

    print(f"WROTE: {texts_out}")
    print(f"WROTE: {twins_out}")
    print(f"WROTE: {html_out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
