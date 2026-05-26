"""
Wave 42 — PPTX Semantic Ingress Runtime.

Purpose:
- perform REAL PPTX ingestion
- extract slide text and notes
- classify slide semantics
- export structured telemetry
- bridge ingestion into convergence-runtime analytics

Scope intentionally narrowed:
- PPTX only
- no Visio/DWG yet
- OCR optional and deferred

Run:
    python src/ingress/pptx_semantic_ingress_runtime_v1.py deck.pptx --out runtime_outputs/pptx_ingress
"""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List

import pandas as pd
from pptx import Presentation


DRAWING_PATTERNS = {
    "P&ID": [r"p&id", r"instrument", r"valve", r"line number"],
    "PFD": [r"process flow", r"mass flow", r"utility flow"],
    "GA": [r"general arrangement", r"layout", r"section"],
    "TOPOLOGY": [r"interface", r"bridge", r"network", r"topology"],
    "COMMISSIONING": [r"commission", r"startup", r"cooldown", r"fill line"],
}


@dataclass
class SlideSemanticRecord:
    slide_number: int
    title: str
    semantic_type: str
    text_length: int
    note_length: int
    contains_image: bool
    contains_table: bool
    extracted_text: str


class PPTXSemanticIngressRuntime:
    """Executable PPTX ingestion runtime."""

    def __init__(self, pptx_path: Path, output_dir: Path):
        self.pptx_path = pptx_path
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def normalize_text(text: str) -> str:
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    @staticmethod
    def classify_semantic(text: str) -> str:
        lowered = text.lower()

        for semantic_type, patterns in DRAWING_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, lowered):
                    return semantic_type

        return "GENERAL"

    @staticmethod
    def extract_notes(slide) -> str:
        try:
            notes_text = []
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text"):
                    notes_text.append(shape.text)
            return "\n".join(notes_text)
        except Exception:
            return ""

    def extract_slide(self, slide, slide_number: int) -> SlideSemanticRecord:
        text_fragments: List[str] = []
        contains_image = False
        contains_table = False

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_fragments.append(shape.text)

            if shape.shape_type == 13:
                contains_image = True

            if getattr(shape, "has_table", False):
                contains_table = True

        slide_text = self.normalize_text("\n".join(text_fragments))
        notes_text = self.normalize_text(self.extract_notes(slide))

        title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_number}"

        semantic = self.classify_semantic(slide_text + " " + notes_text)

        return SlideSemanticRecord(
            slide_number=slide_number,
            title=title,
            semantic_type=semantic,
            text_length=len(slide_text),
            note_length=len(notes_text),
            contains_image=contains_image,
            contains_table=contains_table,
            extracted_text=slide_text[:5000],
        )

    def execute(self) -> Dict[str, str]:
        prs = Presentation(str(self.pptx_path))

        records: List[SlideSemanticRecord] = []

        for idx, slide in enumerate(prs.slides, start=1):
            records.append(self.extract_slide(slide, idx))

        df = pd.DataFrame([asdict(r) for r in records])

        csv_path = self.output_dir / "pptx_semantic_telemetry.csv"
        json_path = self.output_dir / "pptx_semantic_telemetry.json"
        summary_path = self.output_dir / "pptx_summary.json"

        df.to_csv(csv_path, index=False)

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(df.to_dict(orient="records"), f, indent=2)

        summary = {
            "source_pptx": str(self.pptx_path),
            "slide_count": len(records),
            "semantic_distribution": df["semantic_type"].value_counts().to_dict(),
            "contains_images": int(df["contains_image"].sum()),
            "contains_tables": int(df["contains_table"].sum()),
            "mean_slide_text_length": float(df["text_length"].mean()),
        }

        with open(summary_path, "w", encoding="utf-8") as f:
            json.dump(summary, f, indent=2)

        return {
            "csv": str(csv_path),
            "json": str(json_path),
            "summary": str(summary_path),
        }


def main() -> None:
    parser = argparse.ArgumentParser(description="Wave 42 PPTX semantic ingress runtime")
    parser.add_argument("pptx", help="Path to PPTX deck")
    parser.add_argument("--out", default="runtime_outputs/pptx_ingress")

    args = parser.parse_args()

    runtime = PPTXSemanticIngressRuntime(
        pptx_path=Path(args.pptx),
        output_dir=Path(args.out),
    )

    outputs = runtime.execute()
    print(json.dumps(outputs, indent=2))


if __name__ == "__main__":
    main()
