"""
Wave 44 — PPTX Geometry & Semantic Runtime.

Purpose:
- refine PPTX ingress beyond raw text extraction
- extract geometric placement and visual hierarchy
- preserve engineering layout semantics
- export shape topology telemetry

This wave focuses on:
- shape coordinates
- z-order
- grouped shapes
- connector discovery
- title/body classification
- layout-aware semantic telemetry

Run:
    python src/ingress/pptx_geometry_semantic_runtime_v1.py deck.pptx --out runtime_outputs/pptx_geometry
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_TO_MM = 0.000352778


@dataclass
class ShapeTelemetry:
    slide_number: int
    shape_index: int
    shape_name: str
    shape_type: str
    semantic_role: str
    left_mm: float
    top_mm: float
    width_mm: float
    height_mm: float
    z_order: int
    is_grouped: bool
    contains_text: bool
    connector_like: bool
    extracted_text: str


class PPTXGeometrySemanticRuntime:
    """Geometry-aware PPTX ingress runtime."""

    def __init__(self, pptx_path: Path, output_dir: Path):
        self.pptx_path = pptx_path
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def emu_to_mm(value: int) -> float:
        return round(value * EMU_TO_MM, 3)

    @staticmethod
    def semantic_role(shape, text: str) -> str:
        lowered = text.lower()

        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.idx == 0:
                    return "TITLE"
            except Exception:
                pass

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "IMAGE"

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "TABLE"

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "GROUP"

        if "flow" in lowered or "line" in lowered or "pipe" in lowered:
            return "FLOW_OBJECT"

        if len(text.strip()) > 120:
            return "BODY_TEXT"

        return "GENERAL"

    @staticmethod
    def connector_like(shape) -> bool:
        try:
            return shape.shape_type == MSO_SHAPE_TYPE.LINE
        except Exception:
            return False

    def extract_slide_shapes(self, slide, slide_number: int) -> List[ShapeTelemetry]:
        telemetry: List[ShapeTelemetry] = []

        for z_order, shape in enumerate(slide.shapes):
            text = shape.text if hasattr(shape, "text") else ""

            telemetry.append(
                ShapeTelemetry(
                    slide_number=slide_number,
                    shape_index=z_order,
                    shape_name=getattr(shape, "name", f"shape_{z_order}"),
                    shape_type=str(shape.shape_type),
                    semantic_role=self.semantic_role(shape, text),
                    left_mm=self.emu_to_mm(getattr(shape, "left", 0)),
                    top_mm=self.emu_to_mm(getattr(shape, "top", 0)),
                    width_mm=self.emu_to_mm(getattr(shape, "width", 0)),
                    height_mm=self.emu_to_mm(getattr(shape, "height", 0)),
                    z_order=z_order,
                    is_grouped=(shape.shape_type == MSO_SHAPE_TYPE.GROUP),
                    contains_text=bool(text.strip()),
                    connector_like=self.connector_like(shape),
                    extracted_text=text[:2000],
                )
            )

        return telemetry

    def execute(self) -> Dict[str, str]:
        prs = Presentation(str(self.pptx_path))

        all_shapes: List[ShapeTelemetry] = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            all_shapes.extend(self.extract_slide_shapes(slide, slide_number))

        df = pd.DataFrame([asdict(r) for r in all_shapes])

        csv_path = self.output_dir / "pptx_geometry_telemetry.csv"
        json_path = self.output_dir / "pptx_geometry_telemetry.json"
        summary_path = self.output_dir / "pptx_geometry_summary.json"

        df.to_csv(csv_path, index=False)

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(df.to_dict(orient="records"), f, indent=2)

        summary = {
            "slides": int(df["slide_number"].nunique()),
            "shapes": len(df),
            "group_shapes": int(df["is_grouped"].sum()),
            "connector_like_shapes": int(df["connector_like"].sum()),
            "semantic_roles": df["semantic_role"].value_counts().to_dict(),
        }

        with open(summary_path, "w", encoding="utf-8") as f:
            json.dump(summary, f, indent=2)

        return {
            "csv": str(csv_path),
            "json": str(json_path),
            "summary": str(summary_path),
        }


def main() -> None:
    parser = argparse.ArgumentParser(description="Wave 44 PPTX geometry semantic runtime")
    parser.add_argument("pptx", help="Path to PPTX")
    parser.add_argument("--out", default="runtime_outputs/pptx_geometry")

    args = parser.parse_args()

    runtime = PPTXGeometrySemanticRuntime(
        pptx_path=Path(args.pptx),
        output_dir=Path(args.out),
    )

    outputs = runtime.execute()
    print(json.dumps(outputs, indent=2))


if __name__ == "__main__":
    main()
